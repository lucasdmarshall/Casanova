"""Read Office documents to markdown + plain text.

One entry point, :func:`read_office`, auto-detects docx / pptx / xlsx from the
file's contents (they are all ZIP containers, told apart by what's inside) and
dispatches to the right reader. Each reader returns structured **markdown**
(headings, bullets, tables preserved) and a plain-text rendering.

The heavy libraries (python-docx, python-pptx, openpyxl) are imported lazily so
importing hirarareader stays cheap and the tests can target the pieces they need.
"""

from __future__ import annotations

import io
import zipfile
from dataclasses import dataclass, field


class ReaderError(RuntimeError):
    """Raised when a document cannot be detected, opened, or parsed."""


@dataclass
class Document:
    """JSON-ready result of reading an Office document."""

    kind: str
    markdown: str
    text: str
    meta: dict = field(default_factory=dict)
    truncated: bool = False


# --- format detection --------------------------------------------------------

_OOXML_MARKERS = {
    "word/document.xml": "docx",
    "ppt/presentation.xml": "pptx",
    "xl/workbook.xml": "xlsx",
}
_EXT_KIND = {".docx": "docx", ".pptx": "pptx", ".xlsx": "xlsx"}


def detect_kind(data: bytes, filename: str | None = None) -> str:
    """Return 'docx' | 'pptx' | 'xlsx', or raise ReaderError.

    Detection is by content (the OOXML part that only that format contains), so
    a mislabelled extension does not matter. Legacy OLE binaries (.doc/.xls/.ppt)
    are recognized and rejected with a helpful message.
    """
    if data[:4] == b"PK\x03\x04":
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                names = set(zf.namelist())
        except zipfile.BadZipFile:
            names = set()
        for marker, kind in _OOXML_MARKERS.items():
            if marker in names:
                return kind
        # a zip, but not a recognized OOXML — fall through to extension
    elif data[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":
        raise ReaderError(
            "legacy binary Office format (.doc/.xls/.ppt) is not supported — "
            "save as .docx/.xlsx/.pptx and try again"
        )

    if filename:
        for ext, kind in _EXT_KIND.items():
            if filename.lower().endswith(ext):
                return kind
    raise ReaderError("unrecognized document — expected a .docx, .pptx, or .xlsx file")


# --- docx --------------------------------------------------------------------

def _iter_docx_blocks(document):
    """Yield paragraphs and tables in document order."""
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    body = document.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, document)
        elif child.tag == qn("w:tbl"):
            yield Table(child, document)


def _table_to_markdown(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    header = rows[0]
    out = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * width) + " |"]
    for r in rows[1:]:
        out.append("| " + " | ".join(cell.replace("\n", " ").strip() for cell in r) + " |")
    return "\n".join(out)


def read_docx(data: bytes) -> Document:
    try:
        from docx import Document as DocxDocument
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as exc:  # pragma: no cover
        raise ReaderError("python-docx is not installed. pip install python-docx") from exc

    try:
        doc = DocxDocument(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001 — corrupt/unsupported
        raise ReaderError(f"could not open .docx: {exc}") from exc

    md_parts: list[str] = []
    text_parts: list[str] = []
    n_tables = 0
    for block in _iter_docx_blocks(doc):
        if isinstance(block, Paragraph):
            content = block.text.strip()
            if not content:
                continue
            style = (block.style.name if block.style else "") or ""
            if style.startswith("Heading"):
                digits = "".join(c for c in style if c.isdigit())
                level = min(int(digits), 6) if digits else 1
                md_parts.append(f"{'#' * level} {content}")
            elif style.startswith("List"):
                md_parts.append(f"- {content}")
            elif style.lower() == "title":
                md_parts.append(f"# {content}")
            else:
                md_parts.append(content)
            text_parts.append(content)
        elif isinstance(block, Table):
            n_tables += 1
            rows = [[cell.text for cell in row.cells] for row in block.rows]
            md_parts.append(_table_to_markdown(rows))
            for r in rows:
                text_parts.append("\t".join(c.strip() for c in r))

    return Document(
        kind="docx",
        markdown="\n\n".join(md_parts).strip(),
        text="\n".join(text_parts).strip(),
        meta={"paragraphs": len(doc.paragraphs), "tables": n_tables},
    )


# --- pptx --------------------------------------------------------------------

def read_pptx(data: bytes) -> Document:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise ReaderError("python-pptx is not installed. pip install python-pptx") from exc

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise ReaderError(f"could not open .pptx: {exc}") from exc

    md_parts: list[str] = []
    text_parts: list[str] = []
    slides = list(prs.slides)
    for i, slide in enumerate(slides, start=1):
        md_parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_table:
                rows = [
                    [cell.text for cell in row.cells] for row in shape.table.rows
                ]
                md_parts.append(_table_to_markdown(rows))
                for r in rows:
                    text_parts.append("\t".join(c.strip() for c in r))
            elif getattr(shape, "has_text_frame", False):
                content = shape.text_frame.text.strip()
                if content:
                    for line in content.splitlines():
                        line = line.strip()
                        if line:
                            md_parts.append(f"- {line}")
                            text_parts.append(line)
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                md_parts.append(f"> **Notes:** {notes}")
                text_parts.append(f"Notes: {notes}")

    return Document(
        kind="pptx",
        markdown="\n\n".join(md_parts).strip(),
        text="\n".join(text_parts).strip(),
        meta={"slides": len(slides)},
    )


# --- xlsx --------------------------------------------------------------------

def _trim(rows: list[list[str]]) -> list[list[str]]:
    """Drop fully-empty trailing rows and columns."""
    while rows and not any(c.strip() for c in rows[-1]):
        rows.pop()
    if not rows:
        return rows
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    last_col = 0
    for r in rows:
        for i in range(width - 1, -1, -1):
            if r[i].strip():
                last_col = max(last_col, i)
                break
    return [r[: last_col + 1] for r in rows]


def read_xlsx(data: bytes, *, max_rows: int = 2000) -> Document:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover
        raise ReaderError("openpyxl is not installed. pip install openpyxl") from exc

    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:  # noqa: BLE001
        raise ReaderError(f"could not open .xlsx: {exc}") from exc

    md_parts: list[str] = []
    text_parts: list[str] = []
    truncated = False
    sheet_names = wb.sheetnames
    for ws in wb.worksheets:
        md_parts.append(f"## Sheet: {ws.title}")
        rows: list[list[str]] = []
        for r, row in enumerate(ws.iter_rows(values_only=True)):
            if r >= max_rows:
                truncated = True
                break
            rows.append(["" if v is None else str(v) for v in row])
        rows = _trim(rows)
        if rows:
            md_parts.append(_table_to_markdown(rows))
            for r in rows:
                text_parts.append("\t".join(c.strip() for c in r))
    wb.close()

    return Document(
        kind="xlsx",
        markdown="\n\n".join(md_parts).strip(),
        text="\n".join(text_parts).strip(),
        meta={"sheets": len(sheet_names), "sheet_names": sheet_names},
        truncated=truncated,
    )


def read_office(data: bytes, *, filename: str | None = None, max_rows: int = 2000) -> Document:
    """Detect the format and read it. Raises ReaderError on any failure."""
    kind = detect_kind(data, filename)
    if kind == "docx":
        return read_docx(data)
    if kind == "pptx":
        return read_pptx(data)
    return read_xlsx(data, max_rows=max_rows)


__all__ = [
    "Document",
    "ReaderError",
    "detect_kind",
    "read_docx",
    "read_office",
    "read_pptx",
    "read_xlsx",
]
