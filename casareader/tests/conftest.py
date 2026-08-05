"""Shared builders: create tiny real Office files in memory for the tests."""

from __future__ import annotations

import io

import pytest


@pytest.fixture
def docx_bytes() -> bytes:
    from docx import Document

    d = Document()
    d.add_heading("Title Heading", level=1)
    d.add_paragraph("A normal paragraph.")
    d.add_paragraph("Bullet item", style="List Bullet")
    table = d.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "A"
    table.rows[0].cells[1].text = "B"
    table.rows[1].cells[0].text = "1"
    table.rows[1].cells[1].text = "2"
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


@pytest.fixture
def pptx_bytes() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    p = Presentation()
    slide = p.slides.add_slide(p.slide_layouts[5])  # "Title Only"
    slide.shapes.title.text = "Deck Title"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(2))
    box.text_frame.text = "Slide body line"
    slide.notes_slide.notes_text_frame.text = "Some speaker notes"
    buf = io.BytesIO()
    p.save(buf)
    return buf.getvalue()


@pytest.fixture
def xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Name", "Qty"])
    ws.append(["Widget", 5])
    ws.append(["Gadget", 3])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()
